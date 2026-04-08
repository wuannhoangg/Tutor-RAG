from __future__ import annotations

import re
from pathlib import Path
from typing import Any, Dict, List, Tuple
from xml.etree import ElementTree as ET
from zipfile import ZipFile

from ..schemas.document import DocumentMetadata


def _get_attr(obj: Any, *names: str, default: Any = None) -> Any:
    if obj is None:
        return default

    if isinstance(obj, dict):
        for name in names:
            if name in obj:
                return obj[name]
        return default

    for name in names:
        if hasattr(obj, name):
            return getattr(obj, name)
    return default


class PPTXParser:
    """
    Extract text from PPTX files.
    Uses python-pptx when available and falls back to XML parsing.
    """

    def parse(self, file_path: str, source_metadata: DocumentMetadata) -> Tuple[str, List[dict]]:
        path = Path(file_path)

        if path.suffix.lower() != ".pptx":
            raise ValueError("Invalid file type for PPTXParser.")
        if not path.exists():
            raise FileNotFoundError(f"PPTX file not found: {file_path}")

        try:
            slides = self._parse_with_python_pptx(path)
        except Exception:
            slides = self._parse_with_zip_xml(path)

        pieces: List[str] = []
        page_details: List[Dict[str, Any]] = []
        cursor = 0

        for slide_number, slide_text in enumerate(slides, start=1):
            block = f"[Slide {slide_number}]\n{slide_text}".strip()
            start_char = cursor
            pieces.append(block)
            cursor += len(block) + 2

            page_details.append(
                {
                    "page_number": slide_number,
                    "start_char": start_char,
                    "end_char": cursor,
                    "text_excerpt": slide_text[:300],
                }
            )

        raw_text = "\n\n".join(pieces).strip()
        return raw_text, page_details

    def _parse_with_python_pptx(self, path: Path) -> List[str]:
        from pptx import Presentation  # type: ignore

        prs = Presentation(str(path))
        slides: List[str] = []

        for slide in prs.slides:
            fragments: List[str] = []

            for shape in slide.shapes:
                text = getattr(shape, "text", None)
                if text:
                    cleaned = text.strip()
                    if cleaned:
                        fragments.append(cleaned)

            slides.append("\n".join(fragments).strip())

        return slides

    def _parse_with_zip_xml(self, path: Path) -> List[str]:
        slides: List[str] = []
        namespace = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

        with ZipFile(path, "r") as archive:
            slide_paths = [
                name
                for name in archive.namelist()
                if re.match(r"ppt/slides/slide\d+\.xml$", name)
            ]

            slide_paths.sort(
                key=lambda item: int(re.search(r"slide(\d+)\.xml$", item).group(1))
            )

            for slide_path in slide_paths:
                root = ET.fromstring(archive.read(slide_path))
                fragments = []

                for node in root.findall(".//a:t", namespace):
                    if node.text and node.text.strip():
                        fragments.append(node.text.strip())

                slides.append("\n".join(fragments).strip())

        return slides